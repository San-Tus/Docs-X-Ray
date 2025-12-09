import argparse
import json
import re
import csv
from pathlib import Path
from collections import defaultdict
from datetime import datetime

import pdfplumber
from colorama import Fore, Style, init
from tabulate import tabulate
from docx import Document
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
from pptx import Presentation
from striprtf.striprtf import rtf_to_text
from odf import text, teletype
from odf.opendocument import load as odf_load
from odf.table import Table, TableRow, TableCell
from odf.draw import Frame

from modules.html_reporting import generate_html_report

# Initialize colorama for Windows color support
init(autoreset=True)

# Supported file extensions
SUPPORTED_EXTENSIONS = [
    # Documents
    '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.rtf', '.odt', '.ods', '.odp', '.txt',
    # Programming languages
    '.py', '.js', '.ts', '.jsx', '.tsx', '.java', '.cpp', '.c', '.h', '.hpp', '.cs', '.php', '.rb', '.go',
    '.rs', '.swift', '.kt', '.scala', '.r', '.m', '.lua', '.pl', '.sh', '.bash', '.zsh', '.fish', '.ps1',
    # Notebooks
    '.ipynb',
    # Web technologies
    '.html', '.htm', '.css', '.scss', '.sass', '.less', '.vue', '.svelte',
    # Configuration files
    '.json', '.yaml', '.yml', '.toml', '.ini', '.cfg', '.conf', '.config', '.properties', '.env',
    # Markup and data
    '.xml', '.md', '.markdown', '.rst', '.tex', '.csv', '.tsv',
    # Build and package files
    '.gradle', '.maven', '.sbt', '.rake', '.make', '.cmake',
    # Docker and container
    '.dockerfile', '.dockerignore', '.containerfile',
    # CI/CD
    '.gitlab-ci.yml', '.travis.yml', '.circleci',
    # Other config/script files
    '.gitignore', '.gitattributes', '.editorconfig', '.htaccess', '.npmrc', '.babelrc', '.eslintrc',
    '.prettierrc', '.stylelintrc', '.jshintrc', '.ansible', '.terraform', '.tf', '.tfvars'
]


def load_sensitive_words(json_path: Path, case_sensitive: bool = False):
    """
    Load sensitive words and compile them as regex patterns with word boundaries.

    Args:
        json_path: Path to the JSON file containing sensitive words
        case_sensitive: If True, matching will be case-sensitive

    Returns:
        Tuple of (raw data dict, compiled patterns dict)
    """
    with json_path.open("r", encoding="utf-8") as f:
        data = json.load(f)
    compiled = {}
    for category, words in data.items():
        compiled[category] = []
        for word in words:
            # Add word boundaries to match whole words only
            # \b matches word boundaries (between \w and \W)
            pattern = r'\b' + re.escape(word) + r'\b'
            flags = 0 if case_sensitive else re.IGNORECASE
            compiled[category].append(re.compile(pattern, flags))
    return data, compiled


def extract_text_from_docx(file_path: Path):
    """Extract text from DOCX files."""
    try:
        doc = Document(file_path)
        text_pages = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_pages.append(para.text)
        return [" ".join(text_pages)]  # Return as single page
    except Exception as e:
        print(f"{Fore.RED}Error reading DOCX {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_xlsx(file_path: Path):
    """Extract text from XLSX/XLS files."""
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_pages = []
        for sheet in wb.worksheets:
            sheet_text = []
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    sheet_text.append(row_text)
            if sheet_text:
                text_pages.append(" ".join(sheet_text))
        return text_pages
    except Exception as e:
        print(f"{Fore.RED}Error reading Excel {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_pptx(file_path: Path):
    """Extract text from PPTX/PPT files."""
    try:
        prs = Presentation(file_path)
        text_pages = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text_pages.append(" ".join(slide_text))
        return text_pages
    except Exception as e:
        print(f"{Fore.RED}Error reading PowerPoint {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_rtf(file_path: Path):
    """Extract text from RTF files."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            rtf_content = f.read()
            text = rtf_to_text(rtf_content)
            return [text] if text.strip() else []
    except Exception as e:
        print(f"{Fore.RED}Error reading RTF {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_pdf(file_path: Path):
    """Extract text from PDF files."""
    try:
        text_pages = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                if text.strip():
                    text_pages.append(text)
        return text_pages
    except Exception as e:
        print(f"{Fore.RED}Error reading PDF {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_txt(file_path: Path):
    """Extract text from TXT files."""
    try:
        # Try multiple encodings
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252', 'iso-8859-1']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    text = f.read()
                    return [text] if text.strip() else []
            except UnicodeDecodeError:
                continue
        # If all encodings fail, try with errors='ignore'
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
            return [text] if text.strip() else []
    except Exception as e:
        print(f"{Fore.RED}Error reading TXT {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_ipynb(file_path: Path):
    """Extract text from Jupyter Notebook (.ipynb) files."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            notebook = json.load(f)

        text_content = []

        # Extract cells from notebook
        cells = notebook.get('cells', [])
        for cell in cells:
            cell_type = cell.get('cell_type', '')
            source = cell.get('source', [])

            # Convert source to string if it's a list
            if isinstance(source, list):
                cell_text = ''.join(source)
            else:
                cell_text = source

            if cell_text.strip():
                text_content.append(cell_text)

            # Also extract outputs from code cells
            if cell_type == 'code':
                outputs = cell.get('outputs', [])
                for output in outputs:
                    # Extract text/plain outputs
                    if 'text' in output:
                        output_text = output['text']
                        if isinstance(output_text, list):
                            output_text = ''.join(output_text)
                        if output_text.strip():
                            text_content.append(output_text)

                    # Extract data outputs
                    if 'data' in output:
                        data = output['data']
                        if 'text/plain' in data:
                            data_text = data['text/plain']
                            if isinstance(data_text, list):
                                data_text = ''.join(data_text)
                            if data_text.strip():
                                text_content.append(data_text)

        return [' '.join(text_content)] if text_content else []
    except Exception as e:
        print(f"{Fore.RED}Error reading Jupyter Notebook {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_odt(file_path: Path):
    """Extract text from ODT (OpenDocument Text) files."""
    try:
        doc = odf_load(file_path)
        all_paragraphs = doc.getElementsByType(text.P)
        text_content = []
        for para in all_paragraphs:
            para_text = teletype.extractText(para)
            if para_text.strip():
                text_content.append(para_text)
        return [" ".join(text_content)] if text_content else []
    except Exception as e:
        print(f"{Fore.RED}Error reading ODT {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_ods(file_path: Path):
    """Extract text from ODS (OpenDocument Spreadsheet) files."""
    try:
        doc = odf_load(file_path)
        tables = doc.getElementsByType(Table)
        text_pages = []

        for table in tables:
            sheet_text = []
            rows = table.getElementsByType(TableRow)
            for row in rows:
                cells = row.getElementsByType(TableCell)
                row_text = []
                for cell in cells:
                    cell_text = teletype.extractText(cell)
                    if cell_text.strip():
                        row_text.append(cell_text)
                if row_text:
                    sheet_text.append(" ".join(row_text))
            if sheet_text:
                text_pages.append(" ".join(sheet_text))
        return text_pages
    except Exception as e:
        print(f"{Fore.RED}Error reading ODS {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_odp(file_path: Path):
    """Extract text from ODP (OpenDocument Presentation) files."""
    try:
        doc = odf_load(file_path)
        all_paragraphs = doc.getElementsByType(text.P)

        # Group text by frames (slides)
        frames = doc.getElementsByType(Frame)
        text_pages = []

        if frames:
            # Extract text from each frame
            for frame in frames:
                frame_paragraphs = frame.getElementsByType(text.P)
                slide_text = []
                for para in frame_paragraphs:
                    para_text = teletype.extractText(para)
                    if para_text.strip():
                        slide_text.append(para_text)
                if slide_text:
                    text_pages.append(" ".join(slide_text))
        else:
            # Fallback: extract all paragraphs
            slide_text = []
            for para in all_paragraphs:
                para_text = teletype.extractText(para)
                if para_text.strip():
                    slide_text.append(para_text)
            if slide_text:
                text_pages.append(" ".join(slide_text))

        return text_pages
    except Exception as e:
        print(f"{Fore.RED}Error reading ODP {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_file(file_path: Path):
    """Extract text from any supported file format."""
    extension = file_path.suffix.lower()

    # Binary document formats
    if extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif extension in ['.docx', '.doc']:
        return extract_text_from_docx(file_path)
    elif extension in ['.xlsx', '.xls']:
        return extract_text_from_xlsx(file_path)
    elif extension in ['.pptx', '.ppt']:
        return extract_text_from_pptx(file_path)
    elif extension == '.rtf':
        return extract_text_from_rtf(file_path)
    elif extension == '.odt':
        return extract_text_from_odt(file_path)
    elif extension == '.ods':
        return extract_text_from_ods(file_path)
    elif extension == '.odp':
        return extract_text_from_odp(file_path)

    # Jupyter Notebooks (special JSON format)
    elif extension == '.ipynb':
        return extract_text_from_ipynb(file_path)

    # Text-based formats (programming, config, markup files)
    # All code and config files are treated as text files
    elif extension in [
        '.txt', '.py', '.js', '.ts', '.jsx', '.tsx', '.java', '.cpp', '.c', '.h', '.hpp',
        '.cs', '.php', '.rb', '.go', '.rs', '.swift', '.kt', '.scala', '.r', '.m', '.lua',
        '.pl', '.sh', '.bash', '.zsh', '.fish', '.ps1', '.html', '.htm', '.css', '.scss',
        '.sass', '.less', '.vue', '.svelte', '.json', '.yaml', '.yml', '.toml', '.ini',
        '.cfg', '.conf', '.config', '.properties', '.env', '.xml', '.md', '.markdown',
        '.rst', '.tex', '.csv', '.tsv', '.gradle', '.maven', '.sbt', '.rake', '.make',
        '.cmake', '.dockerfile', '.dockerignore', '.containerfile', '.gitlab-ci.yml',
        '.travis.yml', '.circleci', '.gitignore', '.gitattributes', '.editorconfig',
        '.htaccess', '.npmrc', '.babelrc', '.eslintrc', '.prettierrc', '.stylelintrc',
        '.jshintrc', '.ansible', '.terraform', '.tf', '.tfvars'
    ]:
        return extract_text_from_txt(file_path)
    else:
        return []


def get_context(text: str, match_start: int, match_end: int, context_chars: int = 40):
    """Extract surrounding context for a match."""
    start = max(0, match_start - context_chars)
    end = min(len(text), match_end + context_chars)

    before = text[start:match_start]
    matched = text[match_start:match_end]
    after = text[match_end:end]

    # Add ellipsis if truncated
    if start > 0:
        before = "..." + before
    if end < len(text):
        after = after + "..."

    return before, matched, after


def scan_file_for_sensitive_words(file_path: Path, patterns_by_category: dict):
    """Scan any supported file type for sensitive words."""
    results = {}

    # Extract text pages from the file
    text_pages = extract_text_from_file(file_path)

    if not text_pages:
        return results

    for page_index, text in enumerate(text_pages):
        # Normalize whitespace
        text = " ".join(text.split())

        for category, patterns in patterns_by_category.items():
            for pattern in patterns:
                matches = list(pattern.finditer(text))
                if matches:
                    # Extract the actual matched text from the first match
                    # instead of trying to clean the pattern
                    word_clean = matches[0].group(0)

                    if category not in results:
                        results[category] = {}
                    if word_clean not in results[category]:
                        results[category][word_clean] = {
                            'pages': [],
                            'contexts': []
                        }

                    results[category][word_clean]['pages'].append(page_index + 1)
                    # Store first context example
                    if len(results[category][word_clean]['contexts']) == 0:
                        before, matched, after = get_context(text, matches[0].start(), matches[0].end())
                        results[category][word_clean]['contexts'].append({
                            'before': before,
                            'matched': matched,
                            'after': after,
                            'page': page_index + 1
                        })
    return results


def print_results(file_name: str, results: dict):
    """Print colored results with context."""
    print(f"\n{Fore.CYAN}{'='*80}")
    print(f"{Fore.CYAN}📄 File: {Style.BRIGHT}{file_name}")
    print(f"{Fore.CYAN}{'='*80}{Style.RESET_ALL}")

    if not results:
        print(f"{Fore.GREEN}✓ No sensitive terms found.{Style.RESET_ALL}")
        return

    for category, words_dict in results.items():
        print(f"\n{Fore.YELLOW}📂 Category: {Style.BRIGHT}{category}{Style.RESET_ALL}")

        for word, data in words_dict.items():
            pages = data['pages']
            pages_str = ", ".join(map(str, sorted(set(pages))))
            count = len(pages)

            print(f"\n  {Fore.RED}⚠ '{word}'{Style.RESET_ALL} - {Fore.MAGENTA}{count} occurrence(s){Style.RESET_ALL} on page(s): {pages_str}")

            # Show context example
            if data['contexts']:
                ctx = data['contexts'][0]
                print(f"  {Fore.BLUE}Context (page {ctx['page']}): {Style.RESET_ALL}", end="")
                print(f"{ctx['before']}{Fore.RED}{Style.BRIGHT}{ctx['matched']}{Style.RESET_ALL}{ctx['after']}")


def export_statistics_csv(stats: dict, output_path: Path, total_files: int, files_with_hits: int, total_matches: int):
    """Export statistics to CSV format."""
    try:
        with open(output_path, 'w', newline='', encoding='utf-8') as csvfile:
            writer = csv.writer(csvfile)

            # Write header
            writer.writerow(['Category', 'Sensitive Term', 'Total Occurrences'])

            # Write data
            for category, words in sorted(stats.items()):
                for word, count in sorted(words.items(), key=lambda x: x[1], reverse=True):
                    writer.writerow([category, word, count])

            # Write summary
            writer.writerow([])
            writer.writerow(['Summary Statistics'])
            writer.writerow(['Total files scanned', total_files])
            writer.writerow(['Files with hits', files_with_hits])
            writer.writerow(['Total matches found', total_matches])
            writer.writerow(['Unique terms found', sum(len(words) for words in stats.values())])
            writer.writerow(['Categories with findings', len(stats)])

        print(f"{Fore.GREEN}✓ CSV report exported: {Style.BRIGHT}{output_path}{Style.RESET_ALL}")
    except Exception as e:
        print(f"{Fore.RED}✗ Error exporting CSV: {e}{Style.RESET_ALL}")


def export_statistics_xlsx(stats: dict, output_path: Path, total_files: int, files_with_hits: int, total_matches: int):
    """Export statistics to XLSX format with formatting in two sheets."""
    try:
        wb = openpyxl.Workbook()

        # Header formatting
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True)

        # ========== Sheet 1: Occurrences ==========
        ws_occurrences = wb.active
        ws_occurrences.title = "Occurrences"

        # Write headers for Occurrences sheet
        headers = ['Category', 'Sensitive Term', 'Total Occurrences']
        for col, header in enumerate(headers, 1):
            cell = ws_occurrences.cell(row=1, column=col, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal='center')

        # Write data
        row = 2
        for category, words in sorted(stats.items()):
            for word, count in sorted(words.items(), key=lambda x: x[1], reverse=True):
                ws_occurrences.cell(row=row, column=1, value=category)
                ws_occurrences.cell(row=row, column=2, value=word)
                ws_occurrences.cell(row=row, column=3, value=count)
                row += 1

        # Adjust column widths
        ws_occurrences.column_dimensions['A'].width = 25
        ws_occurrences.column_dimensions['B'].width = 30
        ws_occurrences.column_dimensions['C'].width = 20

        # ========== Sheet 2: Statistics ==========
        ws_stats = wb.create_sheet(title="Statistics")

        # Summary formatting
        summary_fill = PatternFill(start_color="E7E6E6", end_color="E7E6E6", fill_type="solid")
        summary_font = Font(bold=True)
        label_font = Font(bold=True)

        # Write summary statistics
        row = 1
        summary_data = [
            ['Metric', 'Value'],
            ['Total files scanned', total_files],
            ['Files with hits', files_with_hits],
            ['Total matches found', total_matches],
            ['Unique terms found', sum(len(words) for words in stats.values())],
            ['Categories with findings', len(stats)]
        ]

        for idx, (label, value) in enumerate(summary_data):
            cell_label = ws_stats.cell(row=row + idx, column=1, value=label)
            cell_value = ws_stats.cell(row=row + idx, column=2, value=value)

            if idx == 0:  # Header row
                cell_label.fill = header_fill
                cell_label.font = header_font
                cell_label.alignment = Alignment(horizontal='center')
                cell_value.fill = header_fill
                cell_value.font = header_font
                cell_value.alignment = Alignment(horizontal='center')
            else:
                cell_label.font = label_font

        # Adjust column widths
        ws_stats.column_dimensions['A'].width = 30
        ws_stats.column_dimensions['B'].width = 20

        wb.save(output_path)
        print(f"{Fore.GREEN}✓ XLSX report exported: {Style.BRIGHT}{output_path}{Style.RESET_ALL}")
    except Exception as e:
        print(f"{Fore.RED}✗ Error exporting XLSX: {e}{Style.RESET_ALL}")


def export_statistics_json(stats: dict, output_path: Path, total_files: int, files_with_hits: int, total_matches: int, file_types: dict):
    """Export statistics to JSON format."""
    try:
        # Prepare data structure
        export_data = {
            "summary": {
                "total_files_scanned": total_files,
                "files_with_hits": files_with_hits,
                "total_matches": total_matches,
                "unique_terms_found": sum(len(words) for words in stats.values()),
                "categories_with_findings": len(stats),
                "scan_timestamp": datetime.now().isoformat()
            },
            "file_types": file_types,
            "findings": {}
        }

        # Convert stats to JSON-friendly format
        for category, words in sorted(stats.items()):
            export_data["findings"][category] = []
            for word, count in sorted(words.items(), key=lambda x: x[1], reverse=True):
                export_data["findings"][category].append({
                    "term": word,
                    "occurrences": count
                })

        with open(output_path, 'w', encoding='utf-8') as jsonfile:
            json.dump(export_data, jsonfile, indent=2, ensure_ascii=False)

        print(f"{Fore.GREEN}✓ JSON report exported: {Style.BRIGHT}{output_path}{Style.RESET_ALL}")
    except Exception as e:
        print(f"{Fore.RED}✗ Error exporting JSON: {e}{Style.RESET_ALL}")


def scan_folder(folder: Path, sensitive_json: Path, sensitivity_list: str, generate_html: bool = True, report_lang: str = "en", recursive: bool = False, output_formats: list = None, output_dir: Path = None, case_sensitive: bool = False):
    _, compiled_patterns = load_sensitive_words(sensitive_json, case_sensitive=case_sensitive)

    # Collect all supported files
    all_files = []
    glob_pattern = "**/*" if recursive else "*"
    for ext in SUPPORTED_EXTENSIONS:
        all_files.extend(folder.glob(f"{glob_pattern}{ext}"))

    all_files = sorted(all_files)

    if not all_files:
        print(f"{Fore.RED}No supported files found in {folder}{Style.RESET_ALL}")
        print(f"{Fore.YELLOW}Supported formats: {', '.join(SUPPORTED_EXTENSIONS)}{Style.RESET_ALL}")
        return

    # Count files by type
    file_types = defaultdict(int)
    for file in all_files:
        file_types[file.suffix.lower()] += 1

    print(f"\n{Fore.GREEN}{Style.BRIGHT}Starting scan of {len(all_files)} file(s)...{Style.RESET_ALL}")
    print(f"{Fore.CYAN}File types found: {dict(file_types)}{Style.RESET_ALL}\n")

    # Statistics tracking
    global_stats = defaultdict(lambda: defaultdict(int))
    total_matches = 0
    files_with_hits = 0
    scan_results = []  # Store results for HTML report

    for idx, file_path in enumerate(all_files, 1):
        # Print progress at the top
        print(f"{Fore.CYAN}[{idx}/{len(all_files)}]{Style.RESET_ALL}", end=" ")

        results = scan_file_for_sensitive_words(file_path, compiled_patterns)

        # Store results for HTML report
        scan_results.append({
            'file_path': str(file_path.absolute()),
            'results': results
        })

        # Update statistics
        if results:
            files_with_hits += 1
        for category, words_dict in results.items():
            for word, data in words_dict.items():
                count = len(data['pages'])
                global_stats[category][word] += count
                total_matches += count

        print_results(file_path.name, results)

    # Print summary statistics
    print_summary_statistics(global_stats, total_matches, len(all_files))

    # Determine output directory
    if output_dir is None:
        output_dir = Path(__file__).resolve().parent
    else:
        output_dir.mkdir(parents=True, exist_ok=True)

    # Generate HTML report
    if generate_html:
        html_output = output_dir / f"scan_report_{sensitivity_list}.html"
        try:
            generate_html_report(
                scan_results=scan_results,
                global_stats=global_stats,
                output_path=html_output,
                sensitivity_list=sensitivity_list,
                total_files=len(all_files),
                files_with_hits=files_with_hits,
                total_matches=total_matches,
                file_types=dict(file_types),
                report_lang=report_lang
            )
            print(f"\n{Fore.GREEN}✓ HTML report generated: {Style.BRIGHT}{html_output}{Style.RESET_ALL}")
        except Exception as e:
            print(f"\n{Fore.RED}✗ Error generating HTML report: {e}{Style.RESET_ALL}")

    # Export statistics in requested formats
    if output_formats and global_stats:
        print(f"\n{Fore.CYAN}{'='*80}")
        print(f"{Fore.CYAN}{Style.BRIGHT}📊 EXPORTING STATISTICS")
        print(f"{Fore.CYAN}{'='*80}{Style.RESET_ALL}\n")

        if 'csv' in output_formats:
            csv_output = output_dir / f"statistics_{sensitivity_list}.csv"
            export_statistics_csv(global_stats, csv_output, len(all_files), files_with_hits, total_matches)

        if 'xlsx' in output_formats:
            xlsx_output = output_dir / f"statistics_{sensitivity_list}.xlsx"
            export_statistics_xlsx(global_stats, xlsx_output, len(all_files), files_with_hits, total_matches)

        if 'json' in output_formats:
            json_output = output_dir / f"statistics_{sensitivity_list}.json"
            export_statistics_json(global_stats, json_output, len(all_files), files_with_hits, total_matches, dict(file_types))


def print_summary_statistics(stats: dict, total_matches: int, total_files: int):
    """Print a summary table of all findings."""
    print(f"\n{Fore.CYAN}{'='*80}")
    print(f"{Fore.CYAN}{Style.BRIGHT}📊 SUMMARY STATISTICS")
    print(f"{Fore.CYAN}{'='*80}{Style.RESET_ALL}\n")

    if not stats:
        print(f"{Fore.GREEN}✓ No sensitive terms found in any files.{Style.RESET_ALL}")
        return

    # Prepare table data
    table_data = []
    for category, words in sorted(stats.items()):
        for word, count in sorted(words.items(), key=lambda x: x[1], reverse=True):
            table_data.append([category, word, count])

    # Print summary table
    headers = [
        f"{Fore.YELLOW}Category{Style.RESET_ALL}",
        f"{Fore.YELLOW}Sensitive Term{Style.RESET_ALL}",
        f"{Fore.YELLOW}Total Occurrences{Style.RESET_ALL}"
    ]

    print(tabulate(table_data, headers=headers, tablefmt="grid"))

    # Overall summary
    print(f"\n{Fore.GREEN}{Style.BRIGHT}Overall Summary:{Style.RESET_ALL}")
    print(f"  • Total files scanned: {Fore.CYAN}{total_files}{Style.RESET_ALL}")
    print(f"  • Total matches found: {Fore.RED}{total_matches}{Style.RESET_ALL}")
    print(f"  • Unique terms found: {Fore.MAGENTA}{sum(len(words) for words in stats.values())}{Style.RESET_ALL}")
    print(f"  • Categories with findings: {Fore.YELLOW}{len(stats)}{Style.RESET_ALL}\n")


def main():
    parser = argparse.ArgumentParser(
        description="Docs X-Ray: Scan documents and code for sensitive words and phrases"
    )
    parser.add_argument(
        "-d",
        "--directory",
        required=True,
        help="Directory containing document files to scan",
    )
    parser.add_argument(
        "-s",
        "--sensitivity-list",
        choices=["cz", "en"],
        default="en",
        help="Sensitive words list to use: 'cz' for Czech or 'en' for English (default: en)",
    )
    parser.add_argument(
        "-r",
        "--recursive",
        action="store_true",
        help="Recursively scan all subdirectories",
    )
    parser.add_argument(
        "--no-html",
        action="store_true",
        help="Disable HTML report generation",
    )
    parser.add_argument(
        "-l",
        "--lang",
        choices=["en", "cz"],
        default="en",
        help="Report language: 'en' for English or 'cz' for Czech (default: en)",
    )
    parser.add_argument(
        "-o",
        "--output-format",
        choices=["csv", "xlsx", "json", "all"],
        help="Export statistics in specified format(s): csv, xlsx, json, or all for all formats",
    )
    parser.add_argument(
        "-O",
        "--output-dir",
        help="Output directory for all reports and statistics (default: script directory)",
    )
    parser.add_argument(
        "-c",
        "--case-sensitive",
        action="store_true",
        help="Enable case-sensitive matching (default: case-insensitive)",
    )
    args = parser.parse_args()

    folder_to_scan = Path(args.directory).expanduser()
    if not folder_to_scan.is_dir():
        print(f"Provided path is not a directory: {folder_to_scan}")
        return

    script_dir = Path(__file__).resolve().parent
    sensitive_words_file = script_dir / f"sensitive_words_{args.sensitivity_list}.json"

    if not sensitive_words_file.is_file():
        print(f"{Fore.RED}Sensitive words file not found: {sensitive_words_file}{Style.RESET_ALL}")
        print(f"{Fore.YELLOW}Please ensure the file exists in the script directory.{Style.RESET_ALL}")
        return

    # Parse output formats
    output_formats = None
    if args.output_format:
        if args.output_format == "all":
            output_formats = ["csv", "xlsx", "json"]
        else:
            output_formats = [args.output_format]

    # Parse output directory
    output_dir = None
    if args.output_dir:
        output_dir = Path(args.output_dir).expanduser()

    print(f"{Fore.GREEN}Using sensitivity list: {Style.BRIGHT}{args.sensitivity_list.upper()}{Style.RESET_ALL} ({sensitive_words_file.name})")
    if args.recursive:
        print(f"{Fore.CYAN}Recursive mode: {Style.BRIGHT}ENABLED{Style.RESET_ALL} - scanning all subdirectories")
    if args.case_sensitive:
        print(f"{Fore.CYAN}Case-sensitive matching: {Style.BRIGHT}ENABLED{Style.RESET_ALL}")
    if output_dir:
        print(f"{Fore.CYAN}Output directory: {Style.BRIGHT}{output_dir}{Style.RESET_ALL}")

    scan_folder(
        folder_to_scan,
        sensitive_words_file,
        args.sensitivity_list,
        generate_html=not args.no_html,
        report_lang=args.lang,
        recursive=args.recursive,
        output_formats=output_formats,
        output_dir=output_dir,
        case_sensitive=args.case_sensitive
    )


if __name__ == "__main__":
    main()
