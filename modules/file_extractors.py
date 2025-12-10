"""
File text extraction functions for various document formats.
"""

import json
from pathlib import Path
from colorama import Fore, Style
from docx import Document
import openpyxl
from pptx import Presentation
from striprtf.striprtf import rtf_to_text
from odf import text, teletype
from odf.opendocument import load as odf_load
from odf.table import Table, TableRow, TableCell
from odf.draw import Frame
import pdfplumber

from .utils import suppress_warnings_and_stderr


def extract_text_from_docx(file_path: Path):
    """Extract text from DOCX files."""
    try:
        doc = Document(file_path)
        text_pages = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_pages.append(para.text)
        return [" ".join(text_pages)]  # Return as single page
    except Exception as e:
        print(f"{Fore.RED}Error reading DOCX {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_xlsx(file_path: Path):
    """Extract text from XLSX/XLS files."""
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_pages = []
        for sheet in wb.worksheets:
            sheet_text = []
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    sheet_text.append(row_text)
            if sheet_text:
                text_pages.append(" ".join(sheet_text))
        return text_pages
    except Exception as e:
        print(f"{Fore.RED}Error reading Excel {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_pptx(file_path: Path):
    """Extract text from PPTX/PPT files."""
    try:
        prs = Presentation(file_path)
        text_pages = []
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text_pages.append(" ".join(slide_text))
        return text_pages
    except Exception as e:
        print(f"{Fore.RED}Error reading PowerPoint {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_rtf(file_path: Path):
    """Extract text from RTF files."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            rtf_content = f.read()
            text = rtf_to_text(rtf_content)
            return [text] if text.strip() else []
    except Exception as e:
        print(f"{Fore.RED}Error reading RTF {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_pdf(file_path: Path):
    """Extract text from PDF files."""
    try:
        text_pages = []
        with suppress_warnings_and_stderr():
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    if text.strip():
                        text_pages.append(text)
        return text_pages
    except Exception as e:
        print(f"{Fore.RED}Error reading PDF {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_txt(file_path: Path):
    """Extract text from TXT files."""
    try:
        # Try multiple encodings
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252', 'iso-8859-1']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    text = f.read()
                    return [text] if text.strip() else []
            except UnicodeDecodeError:
                continue
        # If all encodings fail, try with errors='ignore'
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
            return [text] if text.strip() else []
    except Exception as e:
        print(f"{Fore.RED}Error reading TXT {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_ipynb(file_path: Path):
    """Extract text from Jupyter Notebook (.ipynb) files."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            notebook = json.load(f)

        text_content = []

        # Extract cells from notebook
        cells = notebook.get('cells', [])
        for cell in cells:
            cell_type = cell.get('cell_type', '')
            source = cell.get('source', [])

            # Convert source to string if it's a list
            if isinstance(source, list):
                cell_text = ''.join(source)
            else:
                cell_text = source

            if cell_text.strip():
                text_content.append(cell_text)

            # Also extract outputs from code cells
            if cell_type == 'code':
                outputs = cell.get('outputs', [])
                for output in outputs:
                    # Extract text/plain outputs
                    if 'text' in output:
                        output_text = output['text']
                        if isinstance(output_text, list):
                            output_text = ''.join(output_text)
                        if output_text.strip():
                            text_content.append(output_text)

                    # Extract data outputs
                    if 'data' in output:
                        data = output['data']
                        if 'text/plain' in data:
                            data_text = data['text/plain']
                            if isinstance(data_text, list):
                                data_text = ''.join(data_text)
                            if data_text.strip():
                                text_content.append(data_text)

        return [' '.join(text_content)] if text_content else []
    except Exception as e:
        print(f"{Fore.RED}Error reading Jupyter Notebook {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_odt(file_path: Path):
    """Extract text from ODT (OpenDocument Text) files."""
    try:
        doc = odf_load(file_path)
        all_paragraphs = doc.getElementsByType(text.P)
        text_content = []
        for para in all_paragraphs:
            para_text = teletype.extractText(para)
            if para_text.strip():
                text_content.append(para_text)
        return [" ".join(text_content)] if text_content else []
    except Exception as e:
        print(f"{Fore.RED}Error reading ODT {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_ods(file_path: Path):
    """Extract text from ODS (OpenDocument Spreadsheet) files."""
    try:
        doc = odf_load(file_path)
        tables = doc.getElementsByType(Table)
        text_pages = []

        for table in tables:
            sheet_text = []
            rows = table.getElementsByType(TableRow)
            for row in rows:
                cells = row.getElementsByType(TableCell)
                row_text = []
                for cell in cells:
                    cell_text = teletype.extractText(cell)
                    if cell_text.strip():
                        row_text.append(cell_text)
                if row_text:
                    sheet_text.append(" ".join(row_text))
            if sheet_text:
                text_pages.append(" ".join(sheet_text))
        return text_pages
    except Exception as e:
        print(f"{Fore.RED}Error reading ODS {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_odp(file_path: Path):
    """Extract text from ODP (OpenDocument Presentation) files."""
    try:
        doc = odf_load(file_path)
        all_paragraphs = doc.getElementsByType(text.P)

        # Group text by frames (slides)
        frames = doc.getElementsByType(Frame)
        text_pages = []

        if frames:
            # Extract text from each frame
            for frame in frames:
                frame_paragraphs = frame.getElementsByType(text.P)
                slide_text = []
                for para in frame_paragraphs:
                    para_text = teletype.extractText(para)
                    if para_text.strip():
                        slide_text.append(para_text)
                if slide_text:
                    text_pages.append(" ".join(slide_text))
        else:
            # Fallback: extract all paragraphs
            slide_text = []
            for para in all_paragraphs:
                para_text = teletype.extractText(para)
                if para_text.strip():
                    slide_text.append(para_text)
            if slide_text:
                text_pages.append(" ".join(slide_text))

        return text_pages
    except Exception as e:
        print(f"{Fore.RED}Error reading ODP {file_path.name}: {e}{Style.RESET_ALL}")
        return []


def extract_text_from_file(file_path: Path):
    """Extract text from any supported file format."""
    extension = file_path.suffix.lower()

    # Binary document formats
    if extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif extension in ['.docx', '.doc']:
        return extract_text_from_docx(file_path)
    elif extension in ['.xlsx', '.xls']:
        return extract_text_from_xlsx(file_path)
    elif extension in ['.pptx', '.ppt']:
        return extract_text_from_pptx(file_path)
    elif extension == '.rtf':
        return extract_text_from_rtf(file_path)
    elif extension == '.odt':
        return extract_text_from_odt(file_path)
    elif extension == '.ods':
        return extract_text_from_ods(file_path)
    elif extension == '.odp':
        return extract_text_from_odp(file_path)

    # Jupyter Notebooks (special JSON format)
    elif extension == '.ipynb':
        return extract_text_from_ipynb(file_path)

    # Text-based formats (programming, config, markup files)
    # All code and config files are treated as text files
    elif extension in [
        '.txt', '.py', '.js', '.ts', '.jsx', '.tsx', '.java', '.cpp', '.c', '.h', '.hpp',
        '.cs', '.php', '.rb', '.go', '.rs', '.swift', '.kt', '.scala', '.r', '.m', '.lua',
        '.pl', '.sh', '.bash', '.zsh', '.fish', '.ps1', '.html', '.htm', '.css', '.scss',
        '.sass', '.less', '.vue', '.svelte', '.json', '.yaml', '.yml', '.toml', '.ini',
        '.cfg', '.conf', '.config', '.properties', '.env', '.xml', '.md', '.markdown',
        '.rst', '.tex', '.csv', '.tsv', '.gradle', '.maven', '.sbt', '.rake', '.make',
        '.cmake', '.dockerfile', '.dockerignore', '.containerfile', '.gitlab-ci.yml',
        '.travis.yml', '.circleci', '.gitignore', '.gitattributes', '.editorconfig',
        '.htaccess', '.npmrc', '.babelrc', '.eslintrc', '.prettierrc', '.stylelintrc',
        '.jshintrc', '.ansible', '.terraform', '.tf', '.tfvars'
    ]:
        return extract_text_from_txt(file_path)
    else:
        return []
